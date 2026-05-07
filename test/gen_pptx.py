#!/usr/bin/env python3
"""
Generate PPTX from the Manulife-Sinochem vs AIA HTML presentation.
Uses python-pptx to create a styled, professional slide deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import copy

# ==========================================
# BRAND COLORS
# ==========================================
GREEN = RGBColor(0x00, 0xA7, 0x58)
GREEN_DARK = RGBColor(0x00, 0x7A, 0x3E)
GREEN_SOFT = RGBColor(0xE6, 0xF5, 0xED)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0F, 0x1A, 0x13)
GRAY_SECONDARY = RGBColor(0x5A, 0x6B, 0x60)
GRAY_LIGHT = RGBColor(0xE5, 0xEB, 0xE7)
BG_OFF = RGBColor(0xF7, 0xF9, 0xF8)
RED_ACCENT = RGBColor(0xD9, 0x4F, 0x5C)
RED_SOFT = RGBColor(0xFD, 0xF2, 0xF3)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_solid_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Noto Sans SC', italic=False):
    """Add a text box with consistent styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Add a rich text box and return its text_frame for manual building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, font_size=16, color=BLACK, bold=False,
             alignment=PP_ALIGN.LEFT, font_name='Noto Sans SC',
             space_after=Pt(4), italic=False, space_before=Pt(0)):
    """Add a paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_shape_rect_sharp(slide, left, top, width, height, fill_color, border_color=None):
    """Add a sharp rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width=Inches(2.0)):
    """Add the brand accent green line under headings."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    return shape


def add_slide_number(slide, current, total):
    """Add page index bottom-right: 01 / 10"""
    tf = add_rich_textbox(slide, Inches(11.8), Inches(6.9), Inches(1.3), Inches(0.4))
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{current:02d}"
    run1.font.size = Pt(9)
    run1.font.color.rgb = BLACK
    run1.font.bold = True
    run1.font.name = 'Manrope'
    p.alignment = PP_ALIGN.RIGHT

    run_div = p.add_run()
    run_div.text = " / "
    run_div.font.size = Pt(9)
    run_div.font.color.rgb = GREEN
    run_div.font.name = 'Manrope'

    run2 = p.add_run()
    run2.text = f"{total:02d}"
    run2.font.size = Pt(9)
    run2.font.color.rgb = GRAY_SECONDARY
    run2.font.name = 'Manrope'


def add_logo_watermark(slide):
    """Add logo + text watermark bottom-left."""
    # Green dot as placeholder for logo
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), Inches(6.95), Inches(0.22), Inches(0.22)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN
    dot.line.fill.background()

    tf = add_rich_textbox(slide, Inches(0.8), Inches(6.9), Inches(1.5), Inches(0.35))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "中宏保险"
    run.font.size = Pt(11)
    run.font.color.rgb = GREEN
    run.font.bold = True
    run.font.name = 'Noto Sans SC'


def add_section_tag(slide, text, top=Inches(0.5)):
    """Add the small section tag like '01 / 公司背景'."""
    add_textbox(slide, Inches(0.8), top, Inches(10), Inches(0.35),
                text, font_size=9, color=GREEN, bold=True)


def add_heading(slide, text, top=Inches(0.85)):
    """Add heading with brand accent line."""
    add_textbox(slide, Inches(0.8), top, Inches(10), Inches(0.7),
                text, font_size=30, color=BLACK, bold=True)
    add_accent_line(slide, Inches(0.8), top + Inches(0.7), Inches(2.0))


def add_compare_card(slide, left, top, width, height, title, items,
                     highlight=False, bg_color=None):
    """Add a comparison card with title and bullet items."""
    if bg_color is None:
        bg_color = GREEN_SOFT if highlight else BG_OFF
    border_color = GREEN if highlight else GRAY_LIGHT

    card = add_shape_rect(slide, left, top, width, height, bg_color, border_color)

    # Title area: badge + company name on one line
    badge_label = "中宏" if highlight else "友邦"
    badge_w = Inches(0.7)

    # Badge background shape
    badge_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.25), top + Inches(0.17),
        badge_w, Inches(0.3)
    )
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = GREEN if highlight else GRAY_LIGHT
    badge_shape.line.fill.background()

    # Badge text (centered on shape)
    badge_text = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.17),
        badge_w, Inches(0.3)
    )
    bt = badge_text.text_frame
    bp = bt.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = badge_label
    br.font.size = Pt(9)
    br.font.color.rgb = WHITE if highlight else GRAY_SECONDARY
    br.font.bold = True
    br.font.name = 'Noto Sans SC'

    # Company name next to badge
    title_text_box = slide.shapes.add_textbox(
        left + Inches(1.1), top + Inches(0.15),
        width - Inches(1.4), Inches(0.35)
    )
    tt = title_text_box.text_frame
    tp = tt.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(13)
    tr.font.color.rgb = BLACK
    tr.font.bold = True
    tr.font.name = 'Noto Sans SC'

    # Items
    items_box = add_rich_textbox(
        slide, left + Inches(0.25), top + Inches(0.6),
        width - Inches(0.5), height - Inches(0.75)
    )
    for i, item in enumerate(items):
        add_para(items_box, "• " + item, font_size=12,
                 color=BLACK if highlight else GRAY_SECONDARY,
                 space_after=Pt(8))


def add_feature_item(slide, left, top, width, height, icon_text, title, desc, icon_bg=None):
    """Add a feature row with icon circle, title, and description."""
    if icon_bg is None:
        icon_bg = GREEN_SOFT

    # Icon circle
    icon = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, Inches(0.55), Inches(0.55)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = icon_bg
    icon.line.fill.background()

    # Icon text (letter/symbol centered in circle)
    icon_tb = slide.shapes.add_textbox(left, top + Inches(0.03), Inches(0.55), Inches(0.55))
    it = icon_tb.text_frame
    it.word_wrap = False
    ip = it.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    ir = ip.add_run()
    ir.text = icon_text
    ir.font.size = Pt(18)
    ir.font.color.rgb = GREEN_DARK
    ir.font.bold = True
    ir.font.name = 'Noto Sans SC'

    # Title
    add_textbox(slide, left + Inches(0.7), top + Inches(0.02), width - Inches(0.7), Inches(0.3),
                title, font_size=15, color=BLACK, bold=True)

    # Description
    add_textbox(slide, left + Inches(0.7), top + Inches(0.32), width - Inches(0.7), Inches(0.3),
                desc, font_size=11, color=GRAY_SECONDARY)


def add_reason_card(slide, left, top, width, height, number, title, desc):
    """Add a reason card with green left border."""
    # Card background
    card = add_shape_rect(slide, left, top, width, height, BG_OFF, GRAY_LIGHT)

    # Green left border
    border = add_shape_rect_sharp(
        slide, left, top, Pt(4), height, GREEN
    )
    border.line.fill.background()

    # Number
    add_textbox(slide, left + Inches(0.15), top + Inches(0.08),
                Inches(0.7), Inches(0.5),
                number, font_size=32, color=GREEN, bold=True,
                font_name='Manrope')

    # Title
    add_textbox(slide, left + Inches(0.85), top + Inches(0.08),
                width - Inches(1.0), Inches(0.35),
                title, font_size=14, color=BLACK, bold=True)

    # Description
    add_textbox(slide, left + Inches(0.85), top + Inches(0.4),
                width - Inches(1.0), Inches(0.4),
                desc, font_size=11, color=GRAY_SECONDARY)


# ================================================================
# SLIDE 1: COVER
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_solid_bg(slide, WHITE)

# Tag
add_textbox(slide, Inches(4.2), Inches(2.0), Inches(5), Inches(0.4),
            "专 业 对 比 分 析", font_size=12, color=GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)

# Title - "中宏保险"
add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.0),
            "中宏保险", font_size=52, color=GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)

# "vs 友邦保险"
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
            "vs 友邦保险", font_size=40, color=BLACK, bold=True,
            alignment=PP_ALIGN.CENTER)

# Accent line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(4.4), Inches(1.8), Pt(3)
)
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

# Subtitle
add_textbox(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
            "为什么中宏保险是您更明智的选择", font_size=18, color=GRAY_SECONDARY,
            alignment=PP_ALIGN.CENTER)

# Logo watermark
add_logo_watermark(slide)


# ================================================================
# SLIDE 2: 公司背景与股东结构
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "01 / 公司背景")
add_heading(slide, "公司背景与股东结构")

# Two comparison cards
add_compare_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.3),
                 "中宏保险 Manulife-Sinochem",
                 [
                     "成立于 1996 年，中国首家中外合资寿险公司",
                     "加拿大宏利金融（150+ 年历史）",
                     "中国中化集团（世界500强央企）",
                     "总部：上海",
                 ],
                 highlight=True)

add_compare_card(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.3),
                 "友邦保险 AIA",
                 [
                     "1919年创立于上海，1992年重返中国内地",
                     "AIA集团（外商独资，香港上市）",
                     "单一外资股东，泛亚集团架构",
                     "总部：香港",
                 ],
                 highlight=False)

add_slide_number(slide, 1, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 3: 中外合资，更懂中国
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "02 / 核心优势")
add_heading(slide, "中外合资，更懂中国")

# Feature items
add_feature_item(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(0.8),
                 "◆", "宏利金融",
                 "全球领先的保险专业经验，150+ 年历史，超万亿加元管理资产")

add_feature_item(slide, Inches(0.8), Inches(3.2), Inches(5.3), Inches(0.8),
                 "■", "中化集团",
                 "世界500强央企，深刻理解中国市场和政策环境")

# Key takeaway box
box = add_shape_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.8),
                     GREEN_SOFT, GREEN)

takeaway = add_rich_textbox(slide, Inches(1.1), Inches(4.7), Inches(11.1), Inches(1.5))
add_para(takeaway, '"外资专业 + 本土资源" 的组合', font_size=16,
         color=GREEN_DARK, bold=True, space_after=Pt(8))
add_para(takeaway,
         "使产品设计和服务更贴近中国消费者的实际需求。相比之下，友邦为外商独资，决策链更多依赖香港总部，本土化响应不如中宏灵活。",
         font_size=12, color=GRAY_SECONDARY, space_after=Pt(4))

add_slide_number(slide, 2, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 4: 市场定位与差异化策略
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "03 / 市场定位")
add_heading(slide, "市场定位与差异化策略")

add_compare_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.0),
                 "中宏保险",
                 [
                     '核心战略："健康管理+保险"深度融合',
                     "行业首创健康管理体系",
                     "品牌调性：温暖、陪伴、健康生活方式",
                     "目标客群：中产家庭、注重健康管理人群",
                 ],
                 highlight=True)

add_compare_card(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.0),
                 "友邦保险",
                 [
                     '核心战略："卓越营销员"制度 + 高端客群',
                     "增值服务作为配套，非战略核心",
                     "品牌调性：高端、精英、国际化",
                     "目标客群：高净值人群、企业高管",
                 ],
                 highlight=False)

# Callout pill
pill = add_shape_rect(slide, Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.45),
                      GREEN_SOFT)
pill_tb = add_rich_textbox(slide, Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.45))
p = pill_tb.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "中宏将健康管理嵌入公司战略核心，而非简单增值配套"
r.font.size = Pt(11)
r.font.color.rgb = GREEN_DARK
r.font.bold = True
r.font.name = 'Noto Sans SC'

add_slide_number(slide, 3, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 5: 行业标杆——"保险+健康管理"
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "04 / 健康管理")
add_heading(slide, '行业标杆 —— "保险 + 健康管理"')

# Three feature items in a row
add_feature_item(slide, Inches(0.8), Inches(1.9), Inches(3.5), Inches(1.5),
                 "▶", "MOVE 计划",
                 '将"被动理赔"转变为"主动健康"，运动获保险奖励')

add_feature_item(slide, Inches(4.9), Inches(1.9), Inches(3.5), Inches(1.5),
                 "♥", "健康管家",
                 "健康评估、就医协助、慢病管理全流程服务")

add_feature_item(slide, Inches(9.0), Inches(1.9), Inches(3.5), Inches(1.5),
                 "▤", "互动式保单",
                 "客户健康行为与保单利益直接挂钩")

# Key stat
stat_box = add_rich_textbox(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.5))
p = stat_box.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "领先行业近 10 年"
r.font.size = Pt(40)
r.font.color.rgb = GREEN
r.font.bold = True
r.font.name = 'Manrope'

add_para(stat_box, "中宏是国内最早将健康管理作为核心战略的保险公司",
         font_size=12, color=GRAY_SECONDARY,
         alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 4, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 6: 产品灵活定制，性价比更高
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "05 / 产品特色")
add_heading(slide, "产品灵活定制，性价比更高")

# Left column: 中宏优势
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(2), Inches(0.35),
            "中宏优势", font_size=12, color=GREEN_DARK, bold=True)
# Green pill background
pill = add_shape_rect(slide, Inches(0.7), Inches(1.75), Inches(1.1), Inches(0.35),
                      GREEN_SOFT)

items_zh = [
    ("✓", "灵活定制化 — 保障组合自由搭配", True),
    ("✓", '"多倍保"系列重疾险 — 多种可选责任', True),
    ("✓", '"健无忧"系列 — 健康管理内嵌到产品', True),
    ("✓", "养老规划丰富 — 依托宏利全球经验", True),
    ("✓", "数字化体验佳 — 全流程线上化", True),
]

y = Inches(2.2)
for icon, text, bold_item in items_zh:
    # Check mark
    check = add_textbox(slide, Inches(1.0), y, Inches(0.3), Inches(0.35),
                        icon, font_size=13, color=GREEN, bold=True)
    add_textbox(slide, Inches(1.3), y, Inches(5.0), Inches(0.35),
                text, font_size=13, color=BLACK)
    # Divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.3), y + Inches(0.35), Inches(5.0), Pt(1)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = GRAY_LIGHT
    div.line.fill.background()
    y += Inches(0.55)

# Right column: 友邦特点
add_textbox(slide, Inches(7.0), Inches(1.8), Inches(2), Inches(0.35),
            "友邦特点", font_size=12, color=NAVY, bold=True)
# Navy pill
add_shape_rect(slide, Inches(6.9), Inches(1.75), Inches(1.1), Inches(0.35),
               RGBColor(0xEA, 0xEF, 0xF4))

items_yb = [
    ("·", "产品定价较高，品牌溢价明显", False),
    ("·", "产品设计偏标准化，定制化空间较小", False),
]

y = Inches(2.2)
for icon, text, _ in items_yb:
    add_textbox(slide, Inches(7.2), y, Inches(0.3), Inches(0.35),
                icon, font_size=13, color=GRAY_SECONDARY)
    add_textbox(slide, Inches(7.5), y, Inches(5.0), Inches(0.35),
                text, font_size=13, color=GRAY_SECONDARY)
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.5), y + Inches(0.35), Inches(5.0), Pt(1)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = GRAY_LIGHT
    div.line.fill.background()
    y += Inches(0.55)

add_slide_number(slide, 5, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 7: 服务更人性化，响应更快捷
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "06 / 服务体验")
add_heading(slide, "服务更人性化，响应更快捷")

# 2x2 grid of feature items
add_feature_item(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(1.0),
                 "✓", "小额快赔", "小额理赔案件处理速度快，流程简便")

add_feature_item(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(1.0),
                 "◉", "全流程线上化", "投保 → 保全 → 理赔，全部在线完成")

add_feature_item(slide, Inches(0.8), Inches(3.5), Inches(5.3), Inches(1.0),
                 "◎", "专属顾问", "为每位客户配备专属顾问，长期陪伴")

add_feature_item(slide, Inches(7.0), Inches(3.5), Inches(5.3), Inches(1.0),
                 "◆", "全国覆盖", "分支机构覆盖全国主要省市")

add_slide_number(slide, 6, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 8: 财务稳健，资源聚焦中国市场
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "07 / 财务实力")
add_heading(slide, "财务稳健，资源聚焦中国市场")

add_feature_item(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(1.0),
                 "▲", "全球系统重要性保险公司",
                 "宏利金融入选 G-SII 名单，全球金融体系核心成员")

add_feature_item(slide, Inches(7.0), Inches(2.0), Inches(5.3), Inches(1.0),
                 "◆", "超万亿加元管理资产",
                 "宏利金融管理资产规模超 1 万亿 加元")

add_feature_item(slide, Inches(0.8), Inches(3.5), Inches(5.3), Inches(1.0),
                 "■", "央企信用背书",
                 "中化集团作为央企股东，提供强有力的信用保障")

add_feature_item(slide, Inches(7.0), Inches(3.5), Inches(5.3), Inches(1.0),
                 "●", "资源优先聚焦中宏",
                 "宏利将全球最优质资源、技术、产品创新优先引入中宏")

add_slide_number(slide, 7, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 9: 为什么更应该选择中宏保险？（五大理由）
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_section_tag(slide, "08 / 核心理由")
add_heading(slide, "为什么更应该选择中宏保险？")

# 2x2 + 1 centered
add_reason_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.6),
                "01", "中外合资独特基因",
                "既懂国际先进经验，又懂中国客户需求")

add_reason_card(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(1.6),
                "02", '行业领先的"保险+健康管理"',
                '从"事后赔付"到"全程守护"')

add_reason_card(slide, Inches(0.8), Inches(3.9), Inches(5.5), Inches(1.6),
                "03", "更高的产品性价比",
                "灵活定制，更适合不同家庭需求")

add_reason_card(slide, Inches(7.0), Inches(3.9), Inches(5.5), Inches(1.6),
                "04", "更温暖的服务体验",
                "有温度的保险，专属顾问全程陪伴")

# Centered fifth reason
add_reason_card(slide, Inches(3.9), Inches(5.8), Inches(5.5), Inches(1.2),
                "05", "深耕中国市场的长期承诺",
                "近 30 年专注中国市场，资源投入更为集中")

add_slide_number(slide, 8, 9)
add_logo_watermark(slide)


# ================================================================
# SLIDE 10: 结尾页
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

# MANULIFE-SINOCHEM
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.4),
            "MANULIFE-SINOCHEM", font_size=11,
            color=RGBColor(0x7A, 0x9A, 0x8A), bold=True,
            alignment=PP_ALIGN.CENTER)

# 中宏保险
add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.0),
            "中宏保险", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# Accent line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(3.8), Inches(1.8), Pt(3)
)
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

# Subtitle
add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5),
            "以人为本，健康守护", font_size=20,
            color=RGBColor(0xAA, 0xCC, 0xBB),
            alignment=PP_ALIGN.CENTER)

# CTA
add_textbox(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.5),
            "选择中宏，选择更安心的未来", font_size=22,
            color=GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)

# URL
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.4),
            "www.manulife-sinochem.com.cn", font_size=12,
            color=RGBColor(0x6A, 0x8A, 0x7A),
            alignment=PP_ALIGN.CENTER)


# ================================================================
# SAVE
# ================================================================
output_path = "/Users/coral/presentations/manulife-sinochem-vs-aia.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
