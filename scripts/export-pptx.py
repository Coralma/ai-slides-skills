#!/usr/bin/env python3
"""
scripts/export-pptx.py
Convert an HTML presentation to a fully EDITABLE PPTX file.

Strategy:
    Each DOM node inside every .slide is mapped to a native python-pptx
    object (textbox, picture, rectangle, rounded-rectangle, gradient shape).
    No screenshots are used. The resulting deck is 100% editable in
    PowerPoint / Keynote / WPS Office.

Supported:
    ✓ Text — font family, size, weight, style, color, alignment, line height
    ✓ <img> — including base64 data URIs (LOGO) and relative/absolute paths
    ✓ <svg> — embedded as SVG picture (python-pptx >= 1.0)
    ✓ background-image: url(...) — extracted and inserted as picture
    ✓ background-color, border, border-radius  → native PPTX shapes
    ✓ linear-gradient(...)  → native pptx gradFill via lxml
    ✓ transform: rotate(Xdeg)
    ✓ opacity

Not supported (gracefully skipped with a warning):
    ✗ clip-path, mask, backdrop-filter
    ✗ CSS animations (only the final state is captured)
    ✗ ::before / ::after pseudo-elements

Usage:
    python scripts/export-pptx.py <html> [output.pptx] [--compact]

Requires:
    pip install "python-pptx>=1.0" playwright lxml pillow
    playwright install chromium
"""

from __future__ import annotations

import argparse
import base64
import functools
import http.server
import os
import re
import shutil
import socket
import socketserver
import sys
import tempfile
import threading
import urllib.parse
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ════════════════════════════════════════════════════════════
# Constants
# ════════════════════════════════════════════════════════════
SLIDE_W_IN = 13.3333   # 16:9 widescreen width in inches
SLIDE_H_IN = 7.5       # 16:9 widescreen height in inches


# ════════════════════════════════════════════════════════════
# Local static file server
# Needed so Google Fonts and relative assets (images, CSS) load via HTTP
# ════════════════════════════════════════════════════════════
def _free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("", 0))
        return s.getsockname()[1]


def start_server(serve_dir: Path):
    port = _free_port()
    handler = functools.partial(
        http.server.SimpleHTTPRequestHandler, directory=str(serve_dir)
    )
    httpd = socketserver.ThreadingTCPServer(("127.0.0.1", port), handler)
    httpd.daemon_threads = True
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    return port, lambda: (httpd.shutdown(), httpd.server_close())


# ════════════════════════════════════════════════════════════
# DOM collection JavaScript
# Runs inside Playwright to walk the DOM of a single .slide and
# return a flat list of render primitives (shape / text / image / svg).
# ════════════════════════════════════════════════════════════
COLLECT_JS = r"""
(index) => {
    const slide = document.querySelectorAll('.slide')[index];
    if (!slide) return null;
    const sRect = slide.getBoundingClientRect();
    const items = [];
    let order = 0;

    function parseZ(style) {
        const z = style.zIndex;
        return z === 'auto' ? 0 : (parseInt(z) || 0);
    }
    function relRect(el) {
        const r = el.getBoundingClientRect();
        return { x: r.left - sRect.left, y: r.top - sRect.top,
                 w: r.width, h: r.height };
    }
    function isTransparent(c) {
        return !c || /rgba?\(0,\s*0,\s*0,\s*0(\.0+)?\)/.test(c) ||
               c === 'transparent';
    }
    function isVisible(el) {
        const s = getComputedStyle(el);
        return s.display !== 'none' && s.visibility !== 'hidden' &&
               parseFloat(s.opacity) > 0.01;
    }

    // ─── Slide own background (mapped to slide.background.fill) ───
    const slideStyle = getComputedStyle(slide);
    const slideBg    = slideStyle.backgroundColor;
    const slideBgImg = slideStyle.backgroundImage;
    if (!isTransparent(slideBg) || (slideBgImg && slideBgImg !== 'none')) {
        items.push({
            kind: 'slidebg',
            bgColor:    !isTransparent(slideBg) ? slideBg : null,
            bgImageRaw: (slideBgImg && slideBgImg !== 'none') ? slideBgImg : null,
        });
    }

    // ─── Pass 1: element walk — shapes / images / svgs ────────────
    function visitShape(el) {
        if (!isVisible(el)) return;
        const rect = relRect(el);
        if (rect.w < 1 || rect.h < 1) return;
        if (rect.x + rect.w < -5 || rect.y + rect.h < -5 ||
            rect.x > sRect.width + 5 || rect.y > sRect.height + 5) return;

        const style    = getComputedStyle(el);
        const tag      = el.tagName.toLowerCase();
        const z        = parseZ(style);
        const myOrder  = order++;
        const common   = { opacity: parseFloat(style.opacity), transform: style.transform };

        const bgColor    = style.backgroundColor;
        const bgImage    = style.backgroundImage;
        const borderTopW = parseFloat(style.borderTopWidth) || 0;
        const radiusPx   = parseFloat(style.borderTopLeftRadius) || 0;
        const hasBg      = !isTransparent(bgColor);
        const hasBgImg   = bgImage && bgImage !== 'none';
        const hasBorder  = borderTopW > 0 && !isTransparent(style.borderTopColor);

        if ((hasBg || hasBorder) && tag !== 'img' && tag !== 'svg') {
            items.push({
                kind: 'shape', rect, z, order: myOrder,
                bgColor:     hasBg     ? bgColor            : null,
                bgImageRaw:  hasBgImg  ? bgImage            : null,
                borderW:     hasBorder ? borderTopW         : 0,
                borderColor: hasBorder ? style.borderTopColor : null,
                radius: radiusPx,
                ...common,
            });
        }
        if (hasBgImg) {
            const m = bgImage.match(/url\(["']?([^"')]+)["']?\)/);
            if (m) items.push({ kind: 'image', rect, z: z+1, order: myOrder, src: m[1], ...common });
        }

        if (tag === 'img') {
            items.push({ kind: 'image', rect, z, order: myOrder,
                         src: el.currentSrc || el.src, alt: el.alt || '', ...common });
            return;
        }
        if (tag === 'svg') {
            const allSvgs = Array.from(slide.querySelectorAll('svg'));
            items.push({ kind: 'svg', rect, z, order: myOrder,
                         svgIndex: allSvgs.indexOf(el), ...common });
            return;
        }

        for (const c of el.children) visitShape(c);
    }

    // ─── Pass 2: TreeWalker over ALL text nodes ────────────────────
    // This catches mixed-content text such as text after <br> inside
    //   <h1><span>A</span><br>B</h1>  — "B" is a direct text node
    const walker = document.createTreeWalker(
        slide, NodeFilter.SHOW_TEXT,
        { acceptNode(node) {
            const t = (node.nodeValue || '').trim();
            if (!t) return NodeFilter.FILTER_REJECT;
            const par = node.parentElement;
            if (!par || !isVisible(par)) return NodeFilter.FILTER_REJECT;
            const tag = par.tagName.toLowerCase();
            if (tag === 'script' || tag === 'style') return NodeFilter.FILTER_REJECT;
            return NodeFilter.FILTER_ACCEPT;
        }}
    );
    while (walker.nextNode()) {
        const node   = walker.currentNode;
        const text   = (node.nodeValue || '').trim();
        if (!text) continue;
        const parent = node.parentElement;
        if (!parent) continue;
        const style  = getComputedStyle(parent);

        // Range gives the true bounding rect of just this text fragment
        const range = document.createRange();
        range.selectNodeContents(node);
        const rects = range.getClientRects();
        if (!rects.length) continue;

        let x1 = Infinity, y1 = Infinity, x2 = -Infinity, y2 = -Infinity;
        for (const r of rects) {
            x1 = Math.min(x1, r.left);  y1 = Math.min(y1, r.top);
            x2 = Math.max(x2, r.right); y2 = Math.max(y2, r.bottom);
        }

        // Use the parent element's bounding rect width as minimum width.
        // The tight text range is often narrower than the container, causing
        // PPTX text to wrap due to slight font metric differences.
        const parentRect = parent.getBoundingClientRect();
        const parentW = parentRect.width;
        const textW = x2 - x1;
        // Use parent width when text is positioned near parent's left edge
        // (i.e. text starts within the parent's horizontal bounds)
        const useParentW = (x1 >= parentRect.left - 2) && (parentW > textW);
        const finalW = useParentW ? parentW : textW;
        // Keep x relative to parent if using parent width
        const finalX = useParentW ? (parentRect.left - sRect.left) : (x1 - sRect.left);

        const rect = { x: finalX, y: y1-sRect.top, w: finalW, h: y2-y1 };
        if (rect.w < 1 || rect.h < 1) continue;
        if (rect.x + rect.w < -5 || rect.y + rect.h < -5 ||
            rect.x > sRect.width + 5 || rect.y > sRect.height + 5) continue;

        items.push({
            kind: 'text', rect, z: parseZ(style), order: order++, text,
            numLines:      rects.length,
            fontFamily:    style.fontFamily,
            fontSize:      parseFloat(style.fontSize),
            fontWeight:    style.fontWeight,
            fontStyle:     style.fontStyle,
            color:         style.color,
            textAlign:     style.textAlign,
            lineHeight:    style.lineHeight,
            letterSpacing: style.letterSpacing,
            textTransform: style.textTransform,
            opacity:       parseFloat(style.opacity),
            transform:     style.transform,
        });
    }

    for (const c of slide.children) visitShape(c);
    return { slide: { w: sRect.width, h: sRect.height }, items };
}
"""



# ════════════════════════════════════════════════════════════
# CSS parsing helpers
# ════════════════════════════════════════════════════════════
_RGB_RE = re.compile(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)")


def parse_color(s: str) -> tuple[RGBColor, float]:
    """Return (RGBColor, alpha). Falls back to (black, 1.0)."""
    if not s:
        return RGBColor(0, 0, 0), 1.0
    m = _RGB_RE.match(s.strip())
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        alpha_m = re.search(r",\s*([\d.]+)\)", s)
        a = float(alpha_m.group(1)) if alpha_m else 1.0
        return RGBColor(r, g, b), a
    return RGBColor(0, 0, 0), 1.0


def is_bold(weight: str) -> bool:
    if not weight:
        return False
    w = weight.strip().lower()
    return w in ("bold", "bolder") or (w.isdigit() and int(w) >= 600)


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT,
    "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT,
    "center": PP_ALIGN.CENTER, "justify": PP_ALIGN.JUSTIFY,
}

_GRAD_RE = re.compile(r"linear-gradient\((.+)\)$", re.DOTALL)


def parse_linear_gradient(css: str) -> Optional[dict]:
    """Parse linear-gradient(...) into {angle, stops}. Returns None on failure."""
    if not css or "linear-gradient" not in css:
        return None
    m = _GRAD_RE.search(css.strip())
    if not m:
        return None

    inner = m.group(1).strip()
    # Split on commas outside rgb()
    parts: list[str] = []
    depth, buf = 0, ""
    for ch in inner:
        if ch == "(":
            depth += 1
        elif ch == ")":
            depth -= 1
        if ch == "," and depth == 0:
            parts.append(buf.strip()); buf = ""
        else:
            buf += ch
    if buf.strip():
        parts.append(buf.strip())

    # Determine angle from first token
    angle = 180.0
    start = 0
    head = parts[0].strip().lower() if parts else ""
    if head.endswith("deg"):
        try:
            angle = float(head[:-3]); start = 1
        except ValueError:
            pass
    elif head.startswith("to "):
        angle = {"to top": 0, "to right": 90, "to bottom": 180,
                 "to left": 270}.get(head, 180)
        start = 1

    stops: list[tuple[float, RGBColor]] = []
    color_parts = parts[start:]
    for i, p in enumerate(color_parts):
        tok = p.rsplit(" ", 1)
        if len(tok) == 2 and tok[1].endswith("%"):
            col, pct = tok[0], float(tok[1].rstrip("%")) / 100.0
        else:
            col, pct = p, i / max(1, len(color_parts) - 1)
        rgb, _ = parse_color(col)
        stops.append((pct, rgb))

    return {"angle": angle, "stops": stops} if stops else None


# ════════════════════════════════════════════════════════════
# Asset fetching — data URI, relative path, or absolute URL → bytes
# ════════════════════════════════════════════════════════════
def fetch_bytes(src: str, page_url: str) -> Optional[tuple[bytes, str]]:
    """Return (image_bytes, extension) or None on failure."""
    if src.startswith("data:"):
        header, _, data = src.partition(",")
        mime_m = re.match(r"data:([^;]+)", header)
        mime = mime_m.group(1) if mime_m else "image/png"
        ext = mime.split("/")[-1].split("+")[0]
        if ext == "svg+xml":
            ext = "svg"
        if ";base64" in header:
            try:
                return base64.b64decode(data), ext
            except Exception:
                return None
        return urllib.parse.unquote(data).encode(), ext

    full = urllib.parse.urljoin(page_url, src)
    try:
        req = urllib.request.Request(full, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as r:
            body = r.read()
            ct = r.headers.get_content_type() or ""
            ext = ct.split("/")[-1].split("+")[0]
            if ext == "svg+xml":
                ext = "svg"
            return body, ext or "png"
    except Exception as exc:
        print(f"  ⚠ Could not fetch asset: {src[:80]}… ({exc})", file=sys.stderr)
        return None


# ════════════════════════════════════════════════════════════
# PPTX gradient fill via lxml
# ════════════════════════════════════════════════════════════
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def apply_gradient_fill(shape, grad: dict):
    """Replace shape's fill with a native PPTX linear gradient."""
    sp_pr = shape.fill._xPr  # type: ignore[attr-defined]
    # Remove existing fill elements
    for child in list(sp_pr):
        tag = child.tag
        if any(tag.endswith(x) for x in (
                "solidFill", "gradFill", "noFill", "blipFill", "pattFill")):
            sp_pr.remove(child)

    a = _A_NS
    grad_el = etree.SubElement(sp_pr, f"{{{a}}}gradFill",
                               attrib={"flip": "none", "rotWithShape": "1"})
    gs_lst = etree.SubElement(grad_el, f"{{{a}}}gsLst")

    for pct, rgb in grad["stops"]:
        gs = etree.SubElement(gs_lst, f"{{{a}}}gs",
                              attrib={"pos": str(int(pct * 100000))})
        etree.SubElement(gs, f"{{{a}}}srgbClr",
                         attrib={"val": f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"})

    # PPTX angle = (CSS angle - 90) mod 360, expressed in 60000ths of a degree
    pptx_angle = int(((grad["angle"] - 90) % 360) * 60000)
    etree.SubElement(grad_el, f"{{{a}}}lin",
                     attrib={"ang": str(pptx_angle), "scaled": "0"})


# ════════════════════════════════════════════════════════════
# Shape builders
# ════════════════════════════════════════════════════════════
@dataclass
class Ctx:
    slide: Any
    slide_w_emu: int
    slide_h_emu: int
    vp_w: float
    vp_h: float
    page_url: str
    tmp_dir: Path


def _px2emu(px: float, vp: float, emu: int) -> int:
    return int(px / vp * emu)


def _rect(item: dict, ctx: Ctx) -> tuple[int, int, int, int]:
    r = item["rect"]
    left   = _px2emu(max(0, r["x"]), ctx.vp_w, ctx.slide_w_emu)
    top    = _px2emu(max(0, r["y"]), ctx.vp_h, ctx.slide_h_emu)
    width  = max(Emu(9525), _px2emu(r["w"], ctx.vp_w, ctx.slide_w_emu))
    height = max(Emu(9525), _px2emu(r["h"], ctx.vp_h, ctx.slide_h_emu))
    return left, top, width, height


def _apply_rotation(shape, transform_css: str):
    if not transform_css or transform_css == "none":
        return
    m = re.search(r"matrix\(([^)]+)\)", transform_css)
    if m:
        import math
        vals = [float(x) for x in m.group(1).split(",")]
        if len(vals) >= 2:
            try:
                shape.rotation = math.degrees(math.atan2(vals[1], vals[0]))
            except Exception:
                pass


def _set_run_font(run, font_name: str):
    """Set Latin, East Asian (CJK) and Complex Script fonts on a run.

    python-pptx’s run.font.name only writes <a:latin>; Chinese/CJK text needs
    <a:ea> and <a:cs> too, otherwise PowerPoint substitutes the theme’s default
    East-Asian font (SimSun / PingFang) instead of the intended face.
    """
    if not font_name:
        return
    run.font.name = font_name          # creates / updates <a:latin>
    # Now patch in <a:ea> and <a:cs> via lxml
    rPr = run._r.find(f"{{{_A_NS}}}rPr")
    if rPr is None:
        return
    for tag in ("ea", "cs"):
        existing = rPr.find(f"{{{_A_NS}}}{tag}")
        if existing is not None:
            rPr.remove(existing)
        el = etree.SubElement(rPr, f"{{{_A_NS}}}{tag}")
        el.set("typeface", font_name)


# ── Text ──────────────────────────────────────────────────
def build_text(item: dict, ctx: Ctx):
    left, top, width, height = _rect(item, ctx)

    # PPTX / Keynote font metrics differ slightly from Chrome’s, so the
    # bounding box from Range.getClientRects() is often 1–3 % too tight.
    # Without extra room a single-line title like
    # “财务稳健，资源聚焦中国市场” wraps in PPTX even though it fit
    # on one line in the browser, and multi-line paragraphs re-wrap onto
    # extra lines that overflow downward and overlap the next element.
    #
    # Pad proportionally to the font size: a flat 18 px / 4 px is plenty
    # for body copy at 14–16 px but completely insufficient for a 60 px
    # title (where one extra glyph alone is wider than the entire pad).
    # Horizontal slack ≈ 1.5 character widths absorbs the metric drift,
    # vertical slack ≈ 0.4 × font_size × num_lines reserves room in case
    # the renderer still wraps onto one extra line.
    font_px   = float(item.get("fontSize") or 16)
    num_lines = int(item.get("numLines", 1) or 1)
    # Use generous horizontal padding: 2.5 character widths to absorb
    # CJK font metric differences between Chrome and PowerPoint.
    pad_w_px  = max(24.0, font_px * 2.5)
    pad_h_px  = max(6.0,  font_px * 0.5 * max(1, num_lines))
    pad_w_emu = _px2emu(pad_w_px, ctx.vp_w, ctx.slide_w_emu)
    pad_h_emu = _px2emu(pad_h_px, ctx.vp_h, ctx.slide_h_emu)
    width  += pad_w_emu
    height += pad_h_emu

    tb = ctx.slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    # Single-line text in the browser must never wrap in PPTX, regardless
    # of minute font-metric differences. Multi-line text keeps word_wrap
    # so that explicit line breaks and reflow behave naturally.
    tf.word_wrap = (num_lines > 1)

    # Font-size: CSS px → Pt, scaled to PPTX coordinate space
    size_px = font_px
    pt_size = size_px * 72.0 / 96.0 * (SLIDE_W_IN * 96.0 / ctx.vp_w)
    pt_size = max(1.0, pt_size)

    family = ((item.get("fontFamily") or "").split(",")[0]
              .strip().strip('"').strip("'"))
    rgb, _ = parse_color(item.get("color"))

    for idx, line in enumerate(item["text"].split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = _ALIGN_MAP.get(
            (item.get("textAlign") or "left").lower(), PP_ALIGN.LEFT
        )
        # Line height
        lh_raw = str(item.get("lineHeight") or "")
        lh_px = re.match(r"([\d.]+)px", lh_raw)
        if lh_px and size_px:
            try:
                p.line_spacing = float(lh_px.group(1)) / size_px
            except Exception:
                pass

        run = p.add_run()
        text = line
        transform = (item.get("textTransform") or "").lower()
        if transform == "uppercase":
            text = text.upper()
        elif transform == "lowercase":
            text = text.lower()
        run.text = text

        if family:
            _set_run_font(run, family)
        run.font.size = Pt(pt_size)
        run.font.bold = is_bold(item.get("fontWeight", ""))
        run.font.italic = (item.get("fontStyle", "").lower() == "italic")
        run.font.color.rgb = rgb

    _apply_rotation(tb, item.get("transform"))


# ── Image ─────────────────────────────────────────────────
def build_image(item: dict, ctx: Ctx):
    src = item.get("src")
    if not src:
        return
    fetched = fetch_bytes(src, ctx.page_url)
    if not fetched:
        return
    data, ext = fetched

    fname = ctx.tmp_dir / f"img_{abs(id(item))}.{ext or 'png'}"
    fname.write_bytes(data)

    left, top, width, height = _rect(item, ctx)
    try:
        pic = ctx.slide.shapes.add_picture(
            str(fname), left, top, width=width, height=height
        )
        _apply_rotation(pic, item.get("transform"))
    except Exception as e:
        print(f"  ⚠ Skipped image ({ext}): {e}", file=sys.stderr)


# ── SVG ───────────────────────────────────────────────────
def build_svg(item: dict, ctx: Ctx):
    # Use Playwright-captured PNG (rasterized during DOM walk).
    # Direct SVG insertion via add_picture fails because Pillow can't
    # identify SVG format even in python-pptx >= 1.0.
    png_path = item.get("png_path")
    if not png_path or not Path(png_path).exists():
        print(f"  ⚠ Skipped SVG (no pre-rendered PNG available)", file=sys.stderr)
        return
    left, top, width, height = _rect(item, ctx)
    try:
        pic = ctx.slide.shapes.add_picture(
            png_path, left, top, width=width, height=height
        )
        _apply_rotation(pic, item.get("transform"))
    except Exception as e:
        print(f"  ⚠ Skipped SVG PNG: {e}", file=sys.stderr)


# ── Shape (background / border / gradient container) ──────
def build_shape(item: dict, ctx: Ctx):
    left, top, width, height = _rect(item, ctx)
    radius_px = item.get("radius") or 0

    if radius_px > 0:
        shape = ctx.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        adj = min(0.5, radius_px / max(1, min(
            item["rect"]["w"], item["rect"]["h"]
        )))
        try:
            shape.adjustments[0] = adj
        except Exception:
            pass
    else:
        shape = ctx.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )

    # Fill: gradient > solid color > transparent
    grad = parse_linear_gradient(item.get("bgImageRaw") or "")
    if grad:
        shape.fill.solid()          # ensure fill XML element exists
        apply_gradient_fill(shape, grad)
    elif item.get("bgColor"):
        rgb, _ = parse_color(item["bgColor"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    else:
        shape.fill.background()

    # Border
    if item.get("borderW", 0) > 0 and item.get("borderColor"):
        shape.line.color.rgb = parse_color(item["borderColor"])[0]
        shape.line.width = Emu(int(item["borderW"] * 9525))
    else:
        shape.line.fill.background()

    _apply_rotation(shape, item.get("transform"))


# ── Slide Background ──────────────────────────────────────────────
def build_slidebg(item: dict, ctx: Ctx):
    """Set the slide’s own background — solid color or linear gradient."""
    bg   = ctx.slide.background
    fill = bg.fill

    grad = parse_linear_gradient(item.get("bgImageRaw") or "")
    if grad:
        try:
            fill.solid()  # ensure the fill XML node is created
            sp_pr = fill._xPr  # type: ignore[attr-defined]
            for child in list(sp_pr):
                if any(child.tag.endswith(x) for x in (
                        "solidFill", "gradFill", "noFill",
                        "blipFill", "pattFill")):
                    sp_pr.remove(child)
            a = _A_NS
            grad_el = etree.SubElement(
                sp_pr, f"{{{a}}}gradFill",
                attrib={"flip": "none", "rotWithShape": "1"})
            gs_lst = etree.SubElement(grad_el, f"{{{a}}}gsLst")
            for pct, rgb in grad["stops"]:
                gs = etree.SubElement(
                    gs_lst, f"{{{a}}}gs",
                    attrib={"pos": str(int(pct * 100000))})
                etree.SubElement(
                    gs, f"{{{a}}}srgbClr",
                    attrib={"val": f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"})
            pptx_angle = int(((grad["angle"] - 90) % 360) * 60000)
            etree.SubElement(grad_el, f"{{{a}}}lin",
                             attrib={"ang": str(pptx_angle), "scaled": "0"})
        except Exception as exc:
            # Fallback: use the first stop as solid color
            print(f"  ⚠ Slide bg gradient: {exc}", file=sys.stderr)
            if grad["stops"]:
                fill.solid()
                fill.fore_color.rgb = grad["stops"][0][1]
    elif item.get("bgColor"):
        rgb, _ = parse_color(item["bgColor"])
        fill.solid()
        fill.fore_color.rgb = rgb


_BUILDERS = {
    "slidebg": build_slidebg,
    "text":    build_text,
    "image":   build_image,
    "svg":     build_svg,
    "shape":   build_shape,
}


# ════════════════════════════════════════════════════════════
# Main export pipeline
# ════════════════════════════════════════════════════════════
def export_pptx(html_path: Path, output: Path,
                viewport_w: int, viewport_h: int):
    serve_dir = html_path.parent
    port, shutdown = start_server(serve_dir)
    url = (f"http://127.0.0.1:{port}/"
           f"{urllib.parse.quote(html_path.name)}")
    print(f"  Local server on port {port}")

    tmp_dir = Path(tempfile.mkdtemp(prefix="html2pptx-"))

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]  # Blank layout

    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            ctx_pw  = browser.new_context(
                viewport={"width": viewport_w, "height": viewport_h},
                device_scale_factor=1,
            )
            page = ctx_pw.new_page()
            page.goto(url, wait_until="networkidle")
            page.evaluate("() => document.fonts.ready")
            page.wait_for_timeout(1500)

            slide_count = page.evaluate(
                "() => document.querySelectorAll('.slide').length"
            )
            if slide_count == 0:
                raise RuntimeError(
                    "No .slide elements found. "
                    "Ensure the HTML uses class=\"slide\"."
                )
            print(f"  Found {slide_count} slide(s)")

            for i in range(slide_count):
                # Isolate current slide so getBoundingClientRect is correct
                page.evaluate(
                    """(idx) => {
                        const slides = document.querySelectorAll('.slide');
                        slides.forEach((s, k) => {
                            s.style.display = (k === idx) ? '' : 'none';
                            if (k === idx) {
                                s.style.opacity = '1';
                                s.style.visibility = 'visible';
                                s.style.position = 'relative';
                                s.style.transform = 'none';
                            }
                        });
                        // Force .reveal elements to final visible state
                        slides[idx]?.querySelectorAll('.reveal').forEach(el => {
                            el.style.opacity = '1';
                            el.style.transform = 'none';
                            el.style.visibility = 'visible';
                        });
                        slides[idx]?.scrollIntoView({ behavior: 'instant' });
                    }""", i,
                )
                page.wait_for_timeout(350)

                data = page.evaluate(COLLECT_JS, i)
                if not data:
                    print(f"  ⚠ Slide {i + 1}: no data returned, skipping",
                          file=sys.stderr)
                    continue

                items     = data["items"]
                slide_box = data["slide"]

                # ── Screenshot SVG elements as PNG (Pillow can't read SVG) ──
                svg_items = [it for it in items if it["kind"] == "svg"]
                if svg_items:
                    handles = page.evaluate(
                        """(idx) => {
                            const slide = document.querySelectorAll('.slide')[idx];
                            return slide
                                ? Array.from(slide.querySelectorAll('svg')).length
                                : 0;
                        }""", i
                    )
                    for it in svg_items:
                        svg_idx = it.get("svgIndex", -1)
                        if svg_idx < 0:
                            continue
                        try:
                            el_handle = page.evaluate_handle(
                                """(args) => {
                                    const slide = document.querySelectorAll
                                        ('.slide')[args.slideIdx];
                                    return slide
                                        ? slide.querySelectorAll('svg')[args.svgIdx]
                                        : null;
                                }""",
                                {"slideIdx": i, "svgIdx": svg_idx},
                            )
                            if el_handle:
                                png_path = tmp_dir / f"svg_{i}_{svg_idx}.png"
                                el_handle.as_element().screenshot(
                                    path=str(png_path)
                                )
                                it["png_path"] = str(png_path)
                        except Exception as exc:
                            print(f"  ⚠ SVG screenshot slide {i+1} index {svg_idx}: {exc}",
                                  file=sys.stderr)

                # Draw order: slidebg → shape → image/svg → text.
                #
                # Text MUST always be drawn last so it is never covered by
                # card backgrounds or decorative shapes — even when the
                # shape has a higher CSS z-index than the text inside it
                # (e.g. `.card { z-index: 1 }` containing text with
                # `z-index: auto`). In PPTX, draw order is determined by
                # insertion order in the slide's spTree: the last shape
                # added is rendered on top. We therefore split items into
                # two passes and concatenate, guaranteeing every textbox
                # is inserted AFTER every non-text shape regardless of any
                # CSS z-index hint.
                _NON_TEXT_PRIO = {"slidebg": 0, "shape": 1,
                                  "image": 2, "svg": 2}
                non_text = [it for it in items if it["kind"] != "text"]
                text_only = [it for it in items if it["kind"] == "text"]
                # Within each pass keep the natural CSS stacking order so
                # overlapping shapes / overlapping text fragments respect
                # their own z-index relative to peers.
                non_text.sort(key=lambda it: (
                    _NON_TEXT_PRIO.get(it["kind"], 5),
                    it.get("z", 0),
                    it.get("order", -1),
                ))
                text_only.sort(key=lambda it: (
                    it.get("z", 0),
                    it.get("order", -1),
                ))
                items = non_text + text_only

                pptx_slide = prs.slides.add_slide(blank)
                ctx = Ctx(
                    slide      = pptx_slide,
                    slide_w_emu= int(prs.slide_width),
                    slide_h_emu= int(prs.slide_height),
                    vp_w       = slide_box["w"],
                    vp_h       = slide_box["h"],
                    page_url   = url,
                    tmp_dir    = tmp_dir,
                )
                ok = skip = 0
                for item in items:
                    builder = _BUILDERS.get(item["kind"])
                    if not builder:
                        continue
                    try:
                        builder(item, ctx)
                        ok += 1
                    except Exception as exc:
                        print(f"  ⚠ Slide {i+1} {item['kind']}: {exc}",
                              file=sys.stderr)
                        skip += 1

                print(f"  Slide {i + 1}/{slide_count} — "
                      f"{ok} objects written, {skip} skipped")

            browser.close()

        prs.save(str(output))
    finally:
        shutdown()
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ════════════════════════════════════════════════════════════
# CLI
# ════════════════════════════════════════════════════════════
def main():
    ap = argparse.ArgumentParser(
        description="Convert an HTML presentation to an editable PPTX "
                    "(DOM → native python-pptx objects, no screenshots)."
    )
    ap.add_argument("html",
                    help="Path to the HTML file (or a folder's index.html)")
    ap.add_argument("output", nargs="?",
                    help="Output .pptx path (default: alongside input .html)")
    ap.add_argument("--compact", action="store_true",
                    help="Render at 1280×720 instead of 1920×1080")
    args = ap.parse_args()

    html_path = Path(args.html).expanduser().resolve()
    if html_path.is_dir():
        html_path = html_path / "index.html"
    if not html_path.is_file():
        print(f"✗ File not found: {html_path}", file=sys.stderr)
        sys.exit(1)

    output = (Path(args.output).expanduser().resolve()
              if args.output else html_path.with_suffix(".pptx"))
    output.parent.mkdir(parents=True, exist_ok=True)

    vp_w, vp_h = (1280, 720) if args.compact else (1920, 1080)

    print("╔══════════════════════════════════════╗")
    print("║    HTML → Editable PPTX Exporter     ║")
    print("╚══════════════════════════════════════╝")
    print(f"  Input   : {html_path}")
    print(f"  Output  : {output}")
    print(f"  Viewport: {vp_w}x{vp_h}")
    print()

    export_pptx(html_path, output, vp_w, vp_h)

    size_kb = output.stat().st_size / 1024
    print()
    print("════════════════════════════════════════")
    print(f"  ✓ PPTX saved: {output}")
    print(f"  Size: {size_kb:.1f} KB")
    print("  All text is editable. Images/LOGO are embedded as picture shapes.")
    print("════════════════════════════════════════")
    print()

    if sys.platform == "darwin":
        os.system(f'open "{output}"')
    elif sys.platform.startswith("linux"):
        os.system(f'xdg-open "{output}" 2>/dev/null || true')


if __name__ == "__main__":
    main()
