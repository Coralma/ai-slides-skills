#!/usr/bin/env python3
"""
scripts/apply-template.py
Merge PptxGenJS content slides into a company PPTX template.

Strategy
--------
1. Open the company template PPTX.
2. Fill its cover slide (slide 1) placeholders with title / subtitle / date.
3. Copy content slides from a PptxGenJS-generated PPTX, scaling all shape
   coordinates to match the template's canvas size (EMU rescaling).
4. Reorder: cover → content slides → template closing slide.
5. Delete the original template example slides.
6. Save as output PPTX.

Template conventions (中宏PPT模版.pptx)
----------------------------------------
- Slide 1  : cover  — placeholders idx=0 (title), idx=1 (subtitle/author),
                       idx=10 (date)
- Slide 11 : closing — layout 'Closing slide', all visuals in layout
- Slides 2–10 : sample layouts (deleted in output)

Usage
-----
    python3 scripts/apply-template.py \\
        --template  "中宏PPT模版.pptx"  \\
        --input     content.pptx         \\
        --output    final.pptx           \\
        --spec      spec.json            # reads title/subtitle/date/author

    # Or supply text directly:
    python3 scripts/apply-template.py \\
        --template  "中宏PPT模版.pptx"  \\
        --input     content.pptx         \\
        --output    final.pptx           \\
        --title     "演讲标题"           \\
        --subtitle  "副标题"             \\
        --date      "2025年5月"

Options
-------
--skip-first N   Skip the first N slides from --input (default 1, skips the
                 PptxGenJS title slide that is replaced by the template cover).
--skip-last  N   Skip the last  N slides from --input (default 1, skips the
                 PptxGenJS closing slide that is replaced by the template closing).

Requires
--------
    pip install "python-pptx>=1.0" lxml
"""

from __future__ import annotations

import argparse
import copy
import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

# ── XML helpers ──────────────────────────────────────────────────────────────

def rescale_xml(root: etree._Element, sx: float, sy: float) -> etree._Element:
    """
    Scale all shape position / size EMU values in an XML subtree.

    Scaled attributes:
      <a:off  x  y>   — shape / group offsets
      <a:ext  cx cy>  — shape / group extents
      <a:ln   w>      — line width (scaled with sx)

    Intentionally NOT scaled:
      Font sizes  (a:rPr sz  — hundredths of a point, not EMU)
      Shadow blur / offset (kept as absolute pt values)
    """
    for el in root.iter():
        local = el.tag.split("}")[1] if "}" in el.tag else el.tag
        if local == "off":
            if (v := el.get("x")) is not None:
                el.set("x", str(round(int(v) * sx)))
            if (v := el.get("y")) is not None:
                el.set("y", str(round(int(v) * sy)))
        elif local == "ext":
            if (v := el.get("cx")) is not None:
                el.set("cx", str(round(int(v) * sx)))
            if (v := el.get("cy")) is not None:
                el.set("cy", str(round(int(v) * sy)))
        elif local == "ln":
            if (v := el.get("w")) is not None:
                el.set("w", str(round(int(v) * sx)))
    return root


def fill_placeholder(slide, ph_idx: int, text: str) -> bool:
    """
    Set the text of a placeholder (by index), preserving run formatting.
    Returns True if the placeholder was found and updated.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != ph_idx:
            continue
        tf = ph.text_frame
        para = tf.paragraphs[0]
        runs = list(para.runs)
        if runs:
            runs[0].text = text
            for run in runs[1:]:
                run._r.getparent().remove(run._r)
        else:
            run = para.add_run()
            run.text = text
        return True
    return False


# ── Slide copy ───────────────────────────────────────────────────────────────

def copy_content_slides(
    dest_prs: Presentation,
    src_prs:  Presentation,
    src_indices: list[int],
) -> int:
    """
    Append scaled copies of slides (at src_indices) from src_prs to dest_prs.

    Each slide's <p:spTree> and <p:bg> are deep-copied from the source,
    then all EMU coordinates are scaled to fit the destination canvas.

    Returns the number of slides added.
    """
    sx = dest_prs.slide_width.emu  / src_prs.slide_width.emu
    sy = dest_prs.slide_height.emu / src_prs.slide_height.emu

    # Find a Blank layout in the destination to use as the slide base
    blank_layout = None
    for layout in dest_prs.slide_layouts:
        if layout.name.lower() in ("blank", "1_空白", "空白", "1_blank"):
            blank_layout = layout
            break
    if blank_layout is None:
        # Fall back to the last layout
        blank_layout = dest_prs.slide_layouts[-1]

    for idx in src_indices:
        src_slide = src_prs.slides[idx]
        new_slide  = dest_prs.slides.add_slide(blank_layout)

        src_cSld = src_slide._element.find(qn("p:cSld"))
        dst_cSld = new_slide._element.find(qn("p:cSld"))

        if src_cSld is None or dst_cSld is None:
            print(f"  Warning: cSld not found for slide {idx + 1}, skipped")
            continue

        # ── Replace spTree ───────────────────────────────────────────────────
        src_spTree = src_cSld.find(qn("p:spTree"))
        dst_spTree = dst_cSld.find(qn("p:spTree"))
        if src_spTree is not None and dst_spTree is not None:
            new_spTree = copy.deepcopy(src_spTree)
            rescale_xml(new_spTree, sx, sy)
            dst_cSld.replace(dst_spTree, new_spTree)

        # ── Copy background (solid fill defined by PptxGenJS) ────────────────
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            dst_bg = dst_cSld.find(qn("p:bg"))
            new_bg = copy.deepcopy(src_bg)
            if dst_bg is not None:
                dst_cSld.replace(dst_bg, new_bg)
            else:
                # Insert before spTree
                dst_sp = dst_cSld.find(qn("p:spTree"))
                if dst_sp is not None:
                    dst_sp.addprevious(new_bg)
                else:
                    dst_cSld.insert(0, new_bg)

    return len(src_indices)


# ── Slide order management ───────────────────────────────────────────────────

def reorder_slides(prs: Presentation, new_order: list[int]) -> None:
    """
    Reorder all slides in prs according to new_order.
    new_order must be a permutation of range(len(prs.slides)).
    """
    id_list = prs.slides._sldIdLst
    all_ids = list(id_list)
    assert len(new_order) == len(all_ids), (
        f"new_order length {len(new_order)} != slide count {len(all_ids)}"
    )
    for el in list(id_list):
        id_list.remove(el)
    for idx in new_order:
        id_list.append(all_ids[idx])


def delete_last_n_slides(prs: Presentation, n: int) -> None:
    """Remove the last n entries from the slide list (orphans the slide parts)."""
    id_list = prs.slides._sldIdLst
    for _ in range(n):
        entries = list(id_list)
        if entries:
            id_list.remove(entries[-1])


# ── Main ─────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(
        description="Merge PptxGenJS slides into a company PPTX template"
    )
    ap.add_argument("--template",    required=True,  help="Company template .pptx file")
    ap.add_argument("--input",       required=True,  help="PptxGenJS-generated .pptx")
    ap.add_argument("--output",      required=True,  help="Output .pptx path")
    ap.add_argument("--spec",        default="",     help="spec.json (reads title/subtitle/date/author)")
    ap.add_argument("--title",       default="",     help="Cover slide title text")
    ap.add_argument("--subtitle",    default="",     help="Cover slide subtitle text")
    ap.add_argument("--date",        default="",     help="Cover slide date text")
    ap.add_argument("--author",      default="",     help="Cover slide author text")
    ap.add_argument("--skip-first",  type=int, default=1,
                    help="Skip first N slides from --input (default 1)")
    ap.add_argument("--skip-last",   type=int, default=1,
                    help="Skip last N slides from --input (default 1)")
    args = ap.parse_args()

    # ── Read spec (optional) ─────────────────────────────────────────────────
    if args.spec:
        spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
        slides_spec = spec.get("slides", [])
        title_spec  = next((s for s in slides_spec if s.get("type") == "title"), {})
        if not args.title:    args.title    = title_spec.get("title",    spec.get("title",  ""))
        if not args.subtitle: args.subtitle = title_spec.get("subtitle", "")
        if not args.date:     args.date     = title_spec.get("date",     "")
        if not args.author:   args.author   = spec.get("author", "")

    # Normalize title — collapse newlines into spaces for single-line placeholder
    title_text = (
        args.title
        .replace("\\n", " ")
        .replace("\n",  " ")
        .strip()
    )

    # ── Open source files ────────────────────────────────────────────────────
    tpl = Presentation(args.template)
    src = Presentation(args.input)
    n_tpl = len(tpl.slides)
    n_src = len(src.slides)

    print(f"ℹ Template : {n_tpl} slides  "
          f"({tpl.slide_width.inches:.2f}\" × {tpl.slide_height.inches:.2f}\")")
    print(f"ℹ Input    : {n_src} slides  "
          f"({src.slide_width.inches:.2f}\" × {src.slide_height.inches:.2f}\")")

    # ── Fill cover slide placeholders ────────────────────────────────────────
    # Template placeholders: idx=0 title, idx=1 subtitle/author, idx=10 date
    cover = tpl.slides[0]
    ok0  = fill_placeholder(cover,  0, title_text)
    ok1  = fill_placeholder(cover,  1, args.subtitle or args.author)
    ok10 = fill_placeholder(cover, 10, args.date)
    print(f"ℹ Cover    : title={ok0}  subtitle={ok1}  date={ok10}")
    print(f"             '{title_text[:50]}'")

    # ── Determine which input slides to copy ─────────────────────────────────
    sf = args.skip_first
    sl = args.skip_last
    content_indices = list(range(sf, max(sf, n_src - sl)))
    if not content_indices:
        print("⚠ No content slides to copy. Falling back to all input slides.")
        content_indices = list(range(n_src))

    print(f"ℹ Copying  : input slides {[i + 1 for i in content_indices]}")

    # ── Copy content slides (appended to end of tpl) ─────────────────────────
    n_added = copy_content_slides(tpl, src, content_indices)
    print(f"ℹ Added    : {n_added} content slides  (total now {len(tpl.slides)})")

    # ── Reorder ──────────────────────────────────────────────────────────────
    # Current layout in tpl:
    #   [0=cover, 1..n_tpl-2=template_examples, n_tpl-1=closing,
    #    n_tpl..n_tpl+n_added-1=new_content]
    #
    # Desired:
    #   [0=cover, 1..n_added=new_content, n_added+1=closing]
    #   (template examples moved to the end, to be deleted)

    closing_old   = n_tpl - 1
    content_range = list(range(n_tpl, n_tpl + n_added))
    examples_range = list(range(1, n_tpl - 1))          # slides 1..9 in template

    new_order = [0] + content_range + [closing_old] + examples_range
    reorder_slides(tpl, new_order)

    # ── Delete original template example slides (now at the end) ─────────────
    n_examples = len(examples_range)
    delete_last_n_slides(tpl, n_examples)
    print(f"ℹ Deleted  : {n_examples} template example slides")
    print(f"ℹ Final    : {len(tpl.slides)} slides")

    # ── Save ─────────────────────────────────────────────────────────────────
    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    tpl.save(str(out))
    print(f"✓ Saved    : {out.resolve()}")


if __name__ == "__main__":
    main()
